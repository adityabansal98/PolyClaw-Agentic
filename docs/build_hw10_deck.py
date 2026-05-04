"""Generate the HW10 final deck for PolyClaw.

3 slides, exact text from the approved spec. Run with:
    python3 docs/build_hw10_deck.py

Output: docs/HW10-deck.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).parent / "HW10-deck.pptx"

# Color palette: Midnight Executive (navy + ice blue + accent gold)
NAVY = RGBColor(0x1E, 0x27, 0x61)        # primary
ICE = RGBColor(0xCA, 0xDC, 0xFC)         # secondary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x14, 0x1B, 0x3B)     # slightly darker than navy for backgrounds
ACCENT_GOLD = RGBColor(0xFF, 0xC8, 0x4D) # tagline + stats
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE_ARROW = RGBColor(0x60, 0xA5, 0xFA)


def add_dark_background(slide, prs):
    """Fill slide with dark navy background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(
    slide, left, top, width, height, text, *,
    font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
    font_name="Calibri", anchor=MSO_ANCHOR.TOP, italic=False,
):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return tx


def add_bullets(
    slide, left, top, width, height, items, *,
    font_size=13, color=GRAY_LIGHT, bullet_color=ICE,
    line_spacing=1.3, font_name="Calibri",
):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        # Bullet glyph
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.size = Pt(font_size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        bullet.font.name = font_name

        # Item text
        text = p.add_run()
        text.text = item
        text.font.size = Pt(font_size)
        text.font.color.rgb = color
        text.font.name = font_name


def add_outcome_bullets(slide, left, top, width, height, items, glyph, glyph_color):
    """Bullets with a custom glyph (✓, ✗, →) per item."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        p.space_after = Pt(8)

        marker = p.add_run()
        marker.text = f"{glyph}  "
        marker.font.size = Pt(14)
        marker.font.bold = True
        marker.font.color.rgb = glyph_color
        marker.font.name = "Calibri"

        text = p.add_run()
        text.text = item
        text.font.size = Pt(11)
        text.font.color.rgb = GRAY_LIGHT
        text.font.name = "Calibri"


# ────────────────────────────────────────────────────────────────────────────

def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_background(slide, prs)

    # Top-left brand mark
    add_text(
        slide, Inches(0.6), Inches(0.45), Inches(3), Inches(0.4),
        "POLYCLAW", font_size=14, bold=True, color=ACCENT_GOLD,
        font_name="Calibri",
    )

    # Title (huge)
    add_text(
        slide, Inches(0.6), Inches(1.1), Inches(12), Inches(1.3),
        "PolyClaw", font_size=84, bold=True, color=WHITE,
        font_name="Calibri",
    )

    # Tagline (the centerpiece)
    add_text(
        slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.7),
        '"An open platform where AI agents compete on Polymarket"',
        font_size=28, italic=True, color=ACCENT_GOLD,
        font_name="Calibri",
    )

    # Divider line (subtle)
    line = slide.shapes.add_connector(
        1, Inches(0.6), Inches(3.4), Inches(12.7), Inches(3.4)
    )
    line.line.color.rgb = ICE
    line.line.width = Pt(1.0)

    # The Problem
    add_text(
        slide, Inches(0.6), Inches(3.65), Inches(6), Inches(0.4),
        "THE PROBLEM", font_size=12, bold=True, color=ACCENT_GOLD,
    )
    add_text(
        slide, Inches(0.6), Inches(4.0), Inches(6), Inches(2.4),
        "Today, if someone builds an AI agent that wants to bet on prediction "
        "markets, there's nowhere to safely test it, no way to benchmark it "
        "against others, and no shared infrastructure for execution, risk "
        "controls, or performance tracking.",
        font_size=14, color=GRAY_LIGHT,
    )

    # What PolyClaw Is
    add_text(
        slide, Inches(7.0), Inches(3.65), Inches(6), Inches(0.4),
        "WHAT POLYCLAW IS", font_size=12, bold=True, color=ACCENT_GOLD,
    )
    add_text(
        slide, Inches(7.0), Inches(4.0), Inches(6), Inches(2.4),
        "A multi-tenant platform that sits between agents and Polymarket. "
        "Any agent connects through one API and gets backtesting, paper "
        "trading, risk enforcement, and a leaderboard — all out of the box.",
        font_size=14, color=WHITE, bold=False,
    )

    # Stats row at bottom
    stats_y = Inches(6.55)
    stats_h = Inches(0.85)
    stat_w = Inches(4.0)

    for i, (num, label) in enumerate([
        ("30", "agents onboarded"),
        ("11", "API endpoints"),
        ("65", "passing tests"),
    ]):
        x = Inches(0.6 + i * 4.2)
        # Big number
        add_text(
            slide, x, stats_y, stat_w, Inches(0.7),
            num, font_size=44, bold=True, color=ACCENT_GOLD,
            align=PP_ALIGN.CENTER, font_name="Calibri",
        )
        # Label
        add_text(
            slide, x, Inches(7.25), stat_w, Inches(0.4),
            label, font_size=12, color=GRAY_LIGHT,
            align=PP_ALIGN.CENTER,
        )


def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, prs)

    # Title
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
        "What it provides + proof it works",
        font_size=32, bold=True, color=WHITE,
    )

    # ── Architecture diagram (left side) ──
    arch_x = Inches(0.6)
    arch_w = Inches(5.5)
    arch_y_start = Inches(1.4)
    box_h = Inches(1.35)
    gap = Inches(0.4)

    # Section header
    add_text(
        slide, arch_x, Inches(1.05), arch_w, Inches(0.3),
        "ARCHITECTURE", font_size=11, bold=True, color=ACCENT_GOLD,
    )

    # AGENTS layer (top)
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, arch_x, arch_y_start, arch_w, box_h)
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0x2B, 0x36, 0x5C)
    box1.line.color.rgb = ICE
    box1.line.width = Pt(1.0)
    add_text(
        slide, arch_x, arch_y_start + Inches(0.18), arch_w, Inches(0.4),
        "AGENTS  (not ours)", font_size=14, bold=True, color=ICE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, arch_x + Inches(0.2), arch_y_start + Inches(0.6),
        arch_w - Inches(0.4), Inches(0.7),
        "Claude · GPT · Custom Python · MCP clients · LangChain",
        font_size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER,
    )

    # Arrow down
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        arch_x + arch_w/2 - Inches(0.15),
        arch_y_start + box_h + Inches(0.05),
        Inches(0.3), Inches(0.3),
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ACCENT_GOLD
    arrow1.line.fill.background()

    # POLYCLAW layer (middle, accent)
    box2_y = arch_y_start + box_h + gap
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, arch_x, box2_y, arch_w, box_h)
    box2.fill.solid()
    box2.fill.fore_color.rgb = ACCENT_GOLD
    box2.line.color.rgb = ACCENT_GOLD
    box2.line.width = Pt(2.0)
    add_text(
        slide, arch_x, box2_y + Inches(0.18), arch_w, Inches(0.4),
        "POLYCLAW  (what we built)", font_size=14, bold=True, color=DARK_BG,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, arch_x + Inches(0.2), box2_y + Inches(0.6),
        arch_w - Inches(0.4), Inches(0.7),
        "Auth · Risk Gate · Paper Trader · Backtest Queue · Audit Log · Leaderboard",
        font_size=10, color=DARK_BG, align=PP_ALIGN.CENTER, bold=True,
    )

    # Arrow down
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        arch_x + arch_w/2 - Inches(0.15),
        box2_y + box_h + Inches(0.05),
        Inches(0.3), Inches(0.3),
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ACCENT_GOLD
    arrow2.line.fill.background()

    # POLYMARKET layer (bottom)
    box3_y = box2_y + box_h + gap
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, arch_x, box3_y, arch_w, box_h)
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(0x2B, 0x36, 0x5C)
    box3.line.color.rgb = ICE
    box3.line.width = Pt(1.0)
    add_text(
        slide, arch_x, box3_y + Inches(0.18), arch_w, Inches(0.4),
        "POLYMARKET  (not ours)", font_size=14, bold=True, color=ICE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, arch_x + Inches(0.2), box3_y + Inches(0.6),
        arch_w - Inches(0.4), Inches(0.7),
        "Order books · Resolutions · Price history",
        font_size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER,
    )

    # ── Right side: features + results ──
    right_x = Inches(6.7)
    right_w = Inches(6.2)

    add_text(
        slide, right_x, Inches(1.05), right_w, Inches(0.3),
        "WHAT THE PLATFORM PROVIDES", font_size=11, bold=True, color=ACCENT_GOLD,
    )
    add_bullets(
        slide, right_x, Inches(1.4), right_w, Inches(2.6),
        [
            "authenticated API access",
            "backtesting engine with data leakage prevention",
            "paper trading with full audit trail",
            "risk gates that enforce position limits",
            "ranked leaderboard with composite scoring",
        ],
        font_size=12, line_spacing=1.25,
    )

    add_text(
        slide, right_x, Inches(4.25), right_w, Inches(0.3),
        "STRESS-TESTED WITH 30 AGENTS", font_size=11, bold=True, color=ACCENT_GOLD,
    )
    add_bullets(
        slide, right_x, Inches(4.6), right_w, Inches(2.8),
        [
            "risk gate caught 100% of violations",
            "walk-forward detected 3 overfitting agents",
            "kill switch responded in 4.8s",
            "Monte Carlo confidence intervals were accurate for 27/30 agents",
            "best-performing agent hit 1.42 Sharpe",
        ],
        font_size=12, line_spacing=1.25,
        bullet_color=GREEN,
    )


def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
        "What we learned + what's next",
        font_size=32, bold=True, color=WHITE,
    )

    col_w = Inches(4.15)
    col_y = Inches(1.4)
    col_h = Inches(5.8)
    gap_x = Inches(0.15)

    # ── Column 1: What Worked ──
    c1_x = Inches(0.6)
    c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c1_x, col_y, col_w, col_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0x18, 0x36, 0x2B)  # green-tinted dark
    c1.line.color.rgb = GREEN
    c1.line.width = Pt(1.5)

    add_text(
        slide, c1_x + Inches(0.3), col_y + Inches(0.3), col_w - Inches(0.6), Inches(0.4),
        "WHAT WORKED", font_size=14, bold=True, color=GREEN,
    )
    add_outcome_bullets(
        slide,
        c1_x + Inches(0.3), col_y + Inches(0.95),
        col_w - Inches(0.6), col_h - Inches(1.2),
        [
            "Platform handled 30 concurrent agents reliably",
            "risk controls caught every violation with zero false positives",
            "walk-forward validation successfully identified agents gaming in-sample metrics",
        ],
        glyph="✓", glyph_color=GREEN,
    )

    # ── Column 2: What Failed ──
    c2_x = c1_x + col_w + gap_x
    c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c2_x, col_y, col_w, col_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x3B, 0x18, 0x18)  # red-tinted dark
    c2.line.color.rgb = RED
    c2.line.width = Pt(1.5)

    add_text(
        slide, c2_x + Inches(0.3), col_y + Inches(0.3), col_w - Inches(0.6), Inches(0.4),
        "WHAT FAILED", font_size=14, bold=True, color=RED,
    )
    add_outcome_bullets(
        slide,
        c2_x + Inches(0.3), col_y + Inches(0.95),
        col_w - Inches(0.6), col_h - Inches(1.2),
        [
            "Platform slowed down when 30 agents traded at once — had to switch databases mid-project",
            "some agents ran bad strategies and the platform had no way to flag it early",
            "designed for single-threaded backtesting — didn't plan for 30 agents queuing at once",
        ],
        glyph="✗", glyph_color=RED,
    )

    # ── Column 3: Next Steps ──
    c3_x = c2_x + col_w + gap_x
    c3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c3_x, col_y, col_w, col_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(0x18, 0x24, 0x3B)  # blue-tinted dark
    c3.line.color.rgb = BLUE_ARROW
    c3.line.width = Pt(1.5)

    add_text(
        slide, c3_x + Inches(0.3), col_y + Inches(0.3), col_w - Inches(0.6), Inches(0.4),
        "NEXT STEPS", font_size=14, bold=True, color=BLUE_ARROW,
    )
    add_outcome_bullets(
        slide,
        c3_x + Inches(0.3), col_y + Inches(0.95),
        col_w - Inches(0.6), col_h - Inches(1.2),
        [
            "Live Polymarket CLOB integration for real execution",
            "strategy DSL so agents can define logic declaratively",
            "horizontal workers for backtest throughput",
        ],
        glyph="→", glyph_color=BLUE_ARROW,
    )

    # Footer
    add_text(
        slide, Inches(0.6), Inches(7.35), Inches(13), Inches(0.3),
        "github.com/adityabansal98/PolyClaw-Agentic   ·   poly-claw-agentic.vercel.app",
        font_size=11, color=GRAY_MUTED, align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_one(prs)
    slide_two(prs)
    slide_three(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()

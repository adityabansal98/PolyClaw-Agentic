"""Generate a single slide for the Demo Day shared deck (May 7, 2026).

Class format: 35 projects share one deck. Each team gets ONE slide + a 30-second
embedded video. Output: docs/DemoDay-slide.pptx (single slide, 16:9).

Run:
    python3 docs/build_demoday_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).parent / "DemoDay-slide.pptx"

# Palette: dark navy + indigo accent + gold for the tagline
DARK = RGBColor(0x0E, 0x14, 0x29)
NAVY = RGBColor(0x1E, 0x27, 0x61)
INDIGO = RGBColor(0x81, 0x8C, 0xF8)
INDIGO_DEEP = RGBColor(0x63, 0x66, 0xF1)
GOLD = RGBColor(0xFF, 0xC8, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCB, 0xD5, 0xE1)
GRAY_MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)


def text(slide, left, top, w, h, t, **k):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = k.get("anchor", MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = k.get("align", PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = t
    r.font.size = Pt(k.get("size", 14))
    r.font.bold = k.get("bold", False)
    r.font.italic = k.get("italic", False)
    r.font.name = k.get("font", "Calibri")
    r.font.color.rgb = k.get("color", WHITE)
    return tx


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    # ── TOP STRIP: brand + tagline ──────────────────────
    text(slide, Inches(0.5), Inches(0.4), Inches(3), Inches(0.4),
         "POLYCLAW", size=14, bold=True, color=GOLD)

    text(slide, Inches(0.5), Inches(0.85), Inches(8), Inches(0.7),
         "PolyClaw", size=42, bold=True, color=WHITE)

    text(slide, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
         "An open platform where AI agents compete on Polymarket",
         size=18, italic=True, color=GOLD)

    # Live URL — top right
    text(slide, Inches(8.5), Inches(0.5), Inches(4.4), Inches(0.4),
         "poly-claw-agentic.vercel.app", size=11, color=GRAY_MUTED, align=PP_ALIGN.RIGHT)
    text(slide, Inches(8.5), Inches(0.85), Inches(4.4), Inches(0.4),
         "github.com/adityabansal98/PolyClaw-Agentic",
         size=11, color=GRAY_MUTED, align=PP_ALIGN.RIGHT)
    text(slide, Inches(8.5), Inches(1.2), Inches(4.4), Inches(0.4),
         "MIT License · Open source",
         size=10, color=GRAY_MUTED, align=PP_ALIGN.RIGHT)

    # Divider
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(2.25), Inches(12.85), Inches(2.25))
    line.line.color.rgb = INDIGO
    line.line.width = Pt(0.75)

    # ── LEFT COLUMN: What it is ───────────────────────────
    text(slide, Inches(0.5), Inches(2.5), Inches(6), Inches(0.4),
         "WHAT IT IS", size=11, bold=True, color=GOLD)

    text(slide, Inches(0.5), Inches(2.9), Inches(6), Inches(2.4),
         "A multi-tenant platform that sits between agents and Polymarket. "
         "Any agent — Claude, GPT, custom Python, MCP-driven LLMs — connects "
         "through one API and gets backtesting, paper trading, risk enforcement, "
         "and a leaderboard. All out of the box.",
         size=13, color=GRAY)

    # Three-layer mini-diagram
    arch_y = Inches(4.2)
    arch_w = Inches(6)
    box_h = Inches(0.55)

    for i, (label, detail, fill_color, text_color, label_color) in enumerate([
        ("YOUR AGENT", "Claude · GPT · Python · MCP · LangChain", NAVY, WHITE, GRAY),
        ("POLYCLAW", "Auth · Risk · Backtest · Audit · Leaderboard", GOLD, DARK, DARK),
        ("POLYMARKET", "Order books · Resolutions · Prices", NAVY, WHITE, GRAY),
    ]):
        y = arch_y + i * Inches(0.7)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, arch_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.fill.background()

        text(slide, Inches(0.6), y + Inches(0.04), Inches(2.0), Inches(0.3),
             label, size=10, bold=True, color=text_color)
        text(slide, Inches(2.7), y + Inches(0.04), arch_w - Inches(2.3), Inches(0.3),
             detail, size=10, color=label_color)

        if i < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(0.5) + arch_w/2 - Inches(0.1),
                y + box_h + Inches(0.03),
                Inches(0.2), Inches(0.1),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INDIGO
            arrow.line.fill.background()

    # ── RIGHT COLUMN: numbers + outcomes ──────────────────
    text(slide, Inches(7), Inches(2.5), Inches(6), Inches(0.4),
         "STRESS-TESTED WITH 30 AGENTS", size=11, bold=True, color=GOLD)

    # 5 metric cards in a 2x3 grid (5 cards, last cell empty for breathing room)
    metrics = [
        ("100%", "risk gate catch rate", GREEN),
        ("3 / 30", "overfit agents flagged", GOLD),
        ("4.8s", "kill switch response", INDIGO),
        ("27 / 30", "Monte Carlo CI accuracy", GOLD),
        ("1.42", "best agent Sharpe", GREEN),
    ]
    metric_w = Inches(1.95)
    metric_h = Inches(1.25)
    metric_gap = Inches(0.1)

    for i, (num, label, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(7) + col * (metric_w + metric_gap)
        y = Inches(2.95) + row * (metric_h + metric_gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, metric_w, metric_h)
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY
        card.line.color.rgb = INDIGO
        card.line.width = Pt(0.5)

        text(slide, x, y + Inches(0.18), metric_w, Inches(0.55),
             num, size=22, bold=True, color=color, align=PP_ALIGN.CENTER, font="Calibri")
        text(slide, x + Inches(0.1), y + Inches(0.78), metric_w - Inches(0.2), Inches(0.4),
             label, size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # ── BOTTOM STRIP: what's next + try it ────────────────
    text(slide, Inches(0.5), Inches(6.45), Inches(8), Inches(0.4),
         "WHAT'S NEXT", size=11, bold=True, color=GOLD)
    text(slide, Inches(0.5), Inches(6.8), Inches(8), Inches(0.5),
         "Live Polymarket execution · Strategy DSL · Horizontal workers (100+ agents)",
         size=12, color=GRAY)

    # CTA box (right side, bottom)
    cta = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.5), Inches(6.4), Inches(3.4), Inches(0.95),
    )
    cta.fill.solid()
    cta.fill.fore_color.rgb = GOLD
    cta.line.fill.background()
    text(slide, Inches(9.5), Inches(6.55), Inches(3.4), Inches(0.35),
         "TRY IT NOW", size=12, bold=True, color=DARK,
         align=PP_ALIGN.CENTER)
    text(slide, Inches(9.5), Inches(6.92), Inches(3.4), Inches(0.4),
         "poly-claw-agentic.vercel.app",
         size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
